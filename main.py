import sys
from pathlib import Path

from pptx import Presentation
from pptx.slide import Slide, Slides

from screeninfo import get_monitors

import pygame

BASIC_LEVEL_FILE_PATH = "assets/basic-level.pptx"
BACKGROUND_COLOR = (255, 255, 255, 0)
SHAPE_COLOR = (128, 128, 0, 0)
PLAYER_COLOR = (0, 128, 128, 0)

def get_monitor_size() -> tuple[float, float]:
  m = get_monitors()[0]
  return (m.width, m.height)

def get_shapes(filename: str) -> list[pygame.Rect]:
  monitor_size = get_monitor_size()
    
  p = Presentation(filename)
  
  ratio = p.slide_width/monitor_size[0]
  
  first_slide = p.slides[0]
  result = []
  for shape in first_slide.shapes:
    y_pos = shape.top/ratio
    x_pos = shape.left/ratio
    height = shape.height/ratio
    width = shape.width/ratio
    
    rect = pygame.Rect(x_pos, y_pos, width, height)
    result.append(rect)
  return result




def main():
  try:
    shapes = get_shapes(BASIC_LEVEL_FILE_PATH)
    
    pygame.init()
    size = width, height = get_monitor_size()
    screen = pygame.display.set_mode(size)
    
    keep_running = True
    last_time = current_time = pygame.time.get_ticks()
    while keep_running:
      # process inputs
      for event in pygame.event.get():
        if event.type == pygame.QUIT:
          keep_running = False
          break
        elif event.type == pygame.KEYDOWN:
          if event.key == pygame.K_ESCAPE:
            keep_running = False
            break
          else:
            print(f"KEY PRESSED: {event.key}")
      
      # calculate delta time
      current_time = pygame.time.get_ticks()
      delta_time = (current_time - last_time) / 1000.0
      last_time = pygame.time.get_ticks()
      
      # apply updates
      # TODO
      # re-render
      for rect in shapes:
        pygame.draw.rect(screen, SHAPE_COLOR, rect)
      
      pygame.display.update()
    
    pygame.quit()
    sys.exit(0)
  except Exception as ex:
    print(str(ex))
    sys.exit(1)


if __name__ == "__main__":
  main()